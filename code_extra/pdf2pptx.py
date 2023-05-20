## Install relevant packages before running the code
# pip install pdf2image. Windows user may need to do some extra work
# pip install python-pptx

import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import shutil
import sys

# 1. get the relevant info
pdf_name = sys.argv[1]
pages = convert_from_path(pdf_name, 500)
output_folder = os.path.split(pdf_name)[-1].replace(".pdf", "_out")


# 2. make output folder
os.makedirs(output_folder, exist_ok=True)

# 3. read pdf and export pages to output folder
for i, page in enumerate(pages):
    page.save(os.path.join(output_folder, f"page_{str(i).zfill(5)}.png"))

# 4. get the size of an image
image = Image.open(os.path.join(output_folder, f"page_{str(0).zfill(5)}.png"))
width, height = image.width, image.height

# 5. create PPT
# 5.1 Set Slide Size
ppt_pres = Presentation()
ppt_pres.slide_width, ppt_pres.slide_height = Pt(width), Pt(height)
Layout = ppt_pres.slide_layouts[1]
# 5.2 Add all the slides
for j in range(i):
    _ = ppt_pres.slides.add_slide(Layout)
    _.shapes.add_picture(os.path.join(output_folder, f"page_{str(j).zfill(5)}.png"), left=0, top=0)

ppt_pres.save(pdf_name.replace(".pdf",".pptx"))

# 6. Clean up folder
shutil.rmtree(output_folder)

